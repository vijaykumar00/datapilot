"""
report_generator.py — Core report generation and compilation engines.
Supports PDF, DOCX, PPTX, and XLSX bespoke document compilation under thread-safe OOP-style themed Matplotlib rendering.
"""

import os
import logging
import time
from pathlib import Path
from typing import Any, Dict, List
import pandas as pd
import numpy as np

# Configure headless backend for Matplotlib strictly before importing pyplot
import matplotlib
matplotlib.use('Agg')
from matplotlib.figure import Figure
import matplotlib.pyplot as plt

# ReportLab imports for PDF Compiling
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors as rl_colors

# python-docx imports
from docx import Document
from docx.shared import Inches, Pt, RGBColor

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTXPt
from pptx.dml.color import RGBColor as PPTXRGBColor

logger = logging.getLogger("datapilot.report_generator")


def generate_branded_chart(df: pd.DataFrame, x_col: str, y_col: str, chart_type: str, colors: Dict[str, str], save_path: str) -> bool:
    """Generate a themed high-fidelity chart using Object-Oriented thread-safe Matplotlib figure APIs."""
    primary_hex = colors.get("primary", "#6366f1")
    secondary_hex = colors.get("secondary", "#a855f7")

    # Safe validation of columns
    if x_col not in df.columns or y_col not in df.columns:
        logger.warning("Columns '%s' or '%s' do not exist in dataframe. Skipping chart.", x_col, y_col)
        return False

    # Clean missing values
    plot_df = df[[x_col, y_col]].dropna().copy()
    if plot_df.empty:
        return False

    # Use OOP figure creation to ensure multi-threaded concurrency safety
    fig = Figure(figsize=(6.5, 3.8), dpi=150)
    ax = fig.subplots()

    try:
        if chart_type == "bar":
            # Group and take mean
            grouped = plot_df.groupby(x_col)[y_col].mean().head(10)
            ax.bar(grouped.index.astype(str), grouped.values, color=primary_hex, edgecolor=secondary_hex, alpha=0.85, width=0.55)
            ax.set_ylabel(y_col, fontsize=8, color="#555555")
        elif chart_type == "line":
            sorted_df = plot_df.sort_values(by=x_col).head(60)
            ax.plot(sorted_df[x_col].astype(str), sorted_df[y_col], color=primary_hex, marker='o', linestyle='-', linewidth=2.0, markersize=4, alpha=0.9)
            ax.set_ylabel(y_col, fontsize=8, color="#555555")
        elif chart_type == "scatter":
            ax.scatter(plot_df[x_col], plot_df[y_col], color=primary_hex, edgecolors=secondary_hex, alpha=0.7, s=30)
            ax.set_xlabel(x_col, fontsize=8, color="#555555")
            ax.set_ylabel(y_col, fontsize=8, color="#555555")
        else:
            # Default bar
            grouped = plot_df.groupby(x_col)[y_col].mean().head(10)
            ax.bar(grouped.index.astype(str), grouped.values, color=primary_hex, alpha=0.85)
            ax.set_ylabel(y_col, fontsize=8, color="#555555")

        # Visual styling
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#e0e0e0')
        ax.spines['bottom'].set_color('#e0e0e0')
        ax.tick_params(colors='#555555', labelsize=7)
        ax.grid(axis='y', linestyle='--', alpha=0.3)
        ax.set_title(f"{y_col} Breakdown by {x_col}", color='#1f2937', fontsize=9.5, fontweight='bold', pad=12)

        fig.autofmt_xdate(rotation=30)
        fig.tight_layout()

        # Create parent directory if missing
        Path(save_path).parent.mkdir(exist_ok=True, parents=True)
        fig.savefig(save_path, format="png", bbox_inches="tight")
        logger.info("Matplotlib themed figure saved to '%s'", save_path)
        return True
    except Exception as e:
        logger.exception("Failed to render brand Matplotlib chart: %s", e)
        return False
    finally:
        # Secure memory release
        fig.clear()
        plt.close(fig)


def compile_pdf(filepath: str, title: str, date_range: str, narrative: str, kpis: List[Dict[str, Any]], chart_path: str | None, brand_colors: Dict[str, str]) -> None:
    """Compile PDF report using ReportLab flowable elements styled under the designated brand colors."""
    primary_hex = brand_colors.get("primary", "#6366f1")
    primary_rl = rl_colors.HexColor(primary_hex)

    doc = SimpleDocTemplate(filepath, pagesize=letter, rightMargin=45, leftMargin=45, topMargin=45, bottomMargin=45)
    styles = getSampleStyleSheet()

    # Custom typography
    title_style = ParagraphStyle(
        'ReportTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=20,
        leading=24,
        textColor=primary_rl,
        spaceAfter=6
    )
    h2_style = ParagraphStyle(
        'SectionHeading',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=primary_rl,
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True
    )
    body_style = ParagraphStyle(
        'ReportBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=13.5,
        textColor=rl_colors.HexColor("#374151"),
        spaceAfter=8
    )

    story = []

    # Title header block
    story.append(Paragraph(title, title_style))
    if date_range:
        story.append(Paragraph(f"<b>Date Range:</b> {date_range} &nbsp;&nbsp;|&nbsp;&nbsp; <b>Author:</b> DataPilot local AI OS", body_style))
    story.append(Spacer(1, 8))

    # Branded horizontal border line
    line_data = [['']]
    line_table = Table(line_data, colWidths=[520], rowHeights=[2.5])
    line_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), primary_rl),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 0),
        ('TOPPADDING', (0, 0), (-1, -1), 0),
    ]))
    story.append(line_table)
    story.append(Spacer(1, 12))

    # Narrative paragraph loop
    story.append(Paragraph("Overview & Narrative Analysis", h2_style))
    for chunk in narrative.split("\n\n"):
        if chunk.strip():
            story.append(Paragraph(chunk.replace("\n", "<br/>"), body_style))

    # Embedding high-res themed chart
    if chart_path and os.path.exists(chart_path):
        story.append(Spacer(1, 10))
        story.append(Paragraph("Performance Metrics Visualization", h2_style))
        story.append(Image(chart_path, width=440, height=255))
        story.append(Spacer(1, 12))

    # Compiling KPIs Grid
    if kpis:
        story.append(Paragraph("Core Key Performance Indicators (KPIs)", h2_style))
        table_data = [["KPI Metric Name", "Calculated Value", "Severity Badge"]]
        for kpi in kpis:
            table_data.append([
                kpi.get("title", ""),
                kpi.get("metric", ""),
                kpi.get("severity", "info").upper()
            ])

        kpi_table = Table(table_data, colWidths=[260, 130, 130])
        kpi_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), primary_rl),
            ('TEXTCOLOR', (0, 0), (-1, 0), rl_colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BACKGROUND', (0, 1), (-1, -1), rl_colors.HexColor("#fafafa")),
            ('GRID', (0, 0), (-1, -1), 0.5, rl_colors.HexColor("#e5e7eb")),
        ]))
        story.append(kpi_table)

    def draw_footer(canvas, doc):
        canvas.saveState()
        canvas.setFont('Helvetica', 8)
        canvas.setFillColor(rl_colors.HexColor("#6b7280"))
        canvas.drawString(45, 20, "CONFIDENTIAL — DataPilot Automated Document Summary")
        canvas.drawRightString(565, 20, f"Page {doc.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_footer, onLaterPages=draw_footer)


def compile_docx(filepath: str, title: str, date_range: str, narrative: str, kpis: List[Dict[str, Any]], chart_path: str | None, brand_colors: Dict[str, str]) -> None:
    """Compile Word DOCX report using python-docx with custom styled headings and bullet points."""
    primary_hex = brand_colors.get("primary", "#6366f1")
    # Convert HEX to RGB
    r = int(primary_hex.lstrip("#")[0:2], 16)
    g = int(primary_hex.lstrip("#")[2:4], 16)
    b = int(primary_hex.lstrip("#")[4:6], 16)
    brand_rgb = RGBColor(r, g, b)

    doc = Document()

    # Title
    h1 = doc.add_heading(title, level=1)
    for run in h1.runs:
        run.font.color.rgb = brand_rgb
        run.font.size = Pt(20)
        run.font.name = 'Arial'

    if date_range:
        doc.add_paragraph(f"Date Range: {date_range}").runs[0].font.italic = True

    doc.add_paragraph("Generated by DataPilot local AI narratives engine.").runs[0].font.color.rgb = RGBColor(107, 114, 128)

    # Narrative
    h2 = doc.add_heading("Overview & Narrative Analysis", level=2)
    for run in h2.runs:
        run.font.color.rgb = brand_rgb
        run.font.size = Pt(13)

    for chunk in narrative.split("\n\n"):
        if chunk.strip():
            doc.add_paragraph(chunk)

    # Dynamic Chart
    if chart_path and os.path.exists(chart_path):
        h2_c = doc.add_heading("Performance Metrics Visualization", level=2)
        for run in h2_c.runs:
            run.font.color.rgb = brand_rgb
        doc.add_picture(chart_path, width=Inches(5.2))

    # KPIs Grid
    if kpis:
        h2_k = doc.add_heading("Core Key Performance Indicators (KPIs)", level=2)
        for run in h2_k.runs:
            run.font.color.rgb = brand_rgb

        table = doc.add_table(rows=1, cols=3)
        table.style = 'Light Shading Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'KPI Metric Name'
        hdr_cells[1].text = 'Calculated Value'
        hdr_cells[2].text = 'Severity Rating'

        for kpi in kpis:
            row_cells = table.add_row().cells
            row_cells[0].text = kpi.get("title", "")
            row_cells[1].text = kpi.get("metric", "")
            row_cells[2].text = kpi.get("severity", "info").upper()

    doc.save(filepath)


def compile_pptx(filepath: str, title: str, date_range: str, narrative: str, kpis: List[Dict[str, Any]], chart_path: str | None, brand_colors: Dict[str, str]) -> None:
    """Compile PowerPoint PPTX presentation slides under the brand color theme."""
    primary_hex = brand_colors.get("primary", "#6366f1")
    r = int(primary_hex.lstrip("#")[0:2], 16)
    g = int(primary_hex.lstrip("#")[2:4], 16)
    b = int(primary_hex.lstrip("#")[4:6], 16)
    brand_color = PPTXRGBColor(r, g, b)

    prs = Presentation()

    # Slide 1: Cover Layout
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = brand_color
    
    sub = f"Date Range: {date_range}\nGenerated dynamically by DataPilot local AI OS" if date_range else "Generated dynamically by DataPilot local AI OS"
    slide.placeholders[1].text = sub

    # Slide 2: Narrative Analysis
    slide_layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide_layout2)
    slide2.shapes.title.text = "Overview & Narrative Analysis"
    slide2.shapes.title.text_frame.paragraphs[0].font.color.rgb = brand_color

    tf = slide2.placeholders[1].text_frame
    tf.clear()
    for chunk in narrative.split("\n\n")[:3]: # Limit to prevent visual overflow
        if chunk.strip():
            p = tf.add_paragraph()
            p.text = chunk
            p.font.size = PPTXPt(14)
            p.space_after = PPTXPt(8)

    # Slide 3: Embedded Matplotlib themed chart
    if chart_path and os.path.exists(chart_path):
        slide_layout3 = prs.slide_layouts[5] # Title Only
        slide3 = prs.slides.add_slide(slide_layout3)
        slide3.shapes.title.text = "Performance Visualization"
        slide3.shapes.title.text_frame.paragraphs[0].font.color.rgb = brand_color
        slide3.shapes.add_picture(chart_path, PPTXInches(1.5), PPTXInches(1.8), width=PPTXInches(7.0))

    # Slide 4: KPIs Grid Table
    if kpis:
        slide_layout4 = prs.slide_layouts[5]
        slide4 = prs.slides.add_slide(slide_layout4)
        slide4.shapes.title.text = "Key Performance Indicators (KPIs)"
        slide4.shapes.title.text_frame.paragraphs[0].font.color.rgb = brand_color

        rows, cols = len(kpis) + 1, 3
        left, top, width, height = PPTXInches(1.2), PPTXInches(2.0), PPTXInches(7.6), PPTXInches(0.8)
        table = slide4.shapes.add_table(rows, cols, left, top, width, height).table

        table.cell(0, 0).text = "KPI Metric Name"
        table.cell(0, 1).text = "Calculated Value"
        table.cell(0, 2).text = "Severity Badge"

        for idx, kpi in enumerate(kpis):
            table.cell(idx + 1, 0).text = kpi.get("title", "")
            table.cell(idx + 1, 1).text = kpi.get("metric", "")
            table.cell(idx + 1, 2).text = kpi.get("severity", "info").upper()

    prs.save(filepath)


def compile_xlsx(filepath: str, title: str, df: pd.DataFrame, brand_colors: Dict[str, str]) -> None:
    """Compile styled Excel XLSX spreadsheet using openpyxl, themed beautifully with brand headers."""
    from openpyxl.styles import Font, PatternFill, Alignment
    
    writer = pd.ExcelWriter(filepath, engine="openpyxl")
    df.to_excel(writer, sheet_name="Dataset Details", index=False)
    
    workbook = writer.book
    sheet = workbook["Dataset Details"]
    
    primary_hex = brand_colors.get("primary", "6366f1").replace("#", "")
    
    # Styled headers
    header_fill = PatternFill(start_color=primary_hex, end_color=primary_hex, fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    align = Alignment(horizontal="left", vertical="center")
    
    for cell in sheet[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = align
        
    # Auto adjusted column widths
    for col in sheet.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = col[0].column_letter
        sheet.column_dimensions[col_letter].width = max(max_len + 3, 12)
        
    sheet.views.sheetView[0].showGridLines = True
    
    workbook.save(filepath)
    writer.close()
